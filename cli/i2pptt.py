#!/usr/bin/env python3
"""
i2pptt - Image(s) to PowerPoint Tool

Commands:
  - scan:     Scan a directory for images and print metadata
  - group:    Show grouping results based on directory and filename rules
  - generate: Generate a PPTX grouped and laid out by image orientation
"""
from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Dict, List, Optional, Iterable, Tuple

import click
from dataclasses import dataclass, asdict
import configparser
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.util import Inches, Pt
try:
    from pptx.enum.text import PP_ALIGN
except Exception:  # pragma: no cover
    PP_ALIGN = None
try:
    from pptx.enum.shapes import PP_PLACEHOLDER
except Exception:  # pragma: no cover
    PP_PLACEHOLDER = None
import re

# -----------------------------
# Embedded config (from core/config.py)
# -----------------------------
DEFAULTS = {
    "scan": {
        "recursive": "true",
        "supported_formats": "png,jpg,jpeg,gif,webp,bmp",
    },
    "grouping": {
        "strategy": "mixed",
        "filename_separator": "_",
        "max_levels": "3",
    },
    "ppt": {
        "max_images_per_slide": "4",
        "layout_strategy": "auto",
        "language": "zh_CN",
        "size": "16:9",
        "width_in": "",
        "height_in": "",
    },
}


@dataclass
class CliConfig:
    recursive: bool = True
    formats: List[str] = None  # type: ignore[assignment]
    strategy: str = "mixed"
    filename_separator: str = "_"
    max_levels: int = 3
    max_images_per_slide: int = 4
    layout_strategy: str = "auto"
    language: str = "zh_CN"
    ppt_size: str = "16:9"
    ppt_width_in: float | None = None
    ppt_height_in: float | None = None
    landscape_per_slide: int = 3
    portrait_per_slide: int = 4
    square_per_slide: int = 4
    portrait_threshold: float = 0.8
    landscape_threshold: float = 16.0 / 9.0
    title_font_size: int = 28
    content_font_size: int = 10
    title_font_name: str = ""
    content_font_name: str = ""
    title_left_in: float = 0.5
    title_top_in: float = 0.3
    title_width_in: float = 12.4
    title_height_in: float = 1.25
    content_left_in: float = 0.5
    content_top_in: float = 0.8
    content_width_in: float = 13.0
    content_height_in: float = 6.5
    template_path: str = ""

    @staticmethod
    def load(path: Optional[str]) -> "CliConfig":
        cfg = configparser.ConfigParser()
        cfg.read_dict(DEFAULTS)
        if path:
            ini = Path(path).expanduser()
            if not ini.exists():
                raise FileNotFoundError(f"Config not found: {ini}")
            cfg.read(ini, encoding="utf-8")

        recursive = cfg.getboolean("scan", "recursive", fallback=True)
        formats_csv = cfg.get("scan", "supported_formats", fallback=DEFAULTS["scan"]["supported_formats"])
        strategy = cfg.get("grouping", "strategy", fallback="mixed").strip().lower()
        filename_separator = cfg.get("grouping", "filename_separator", fallback="_").strip()
        max_levels = cfg.getint("grouping", "max_levels", fallback=3)
        max_images_per_slide = cfg.getint("ppt", "max_images_per_slide", fallback=4)
        layout_strategy = cfg.get("ppt", "layout_strategy", fallback="auto").strip().lower()
        language = cfg.get("ppt", "language", fallback="zh_CN").strip()
        ppt_size = cfg.get("ppt", "size", fallback="16:9").strip()
        width_in_raw = cfg.get("ppt", "width_in", fallback="").strip()
        height_in_raw = cfg.get("ppt", "height_in", fallback="").strip()
        width_in = float(width_in_raw) if width_in_raw else None
        height_in = float(height_in_raw) if height_in_raw else None
        landscape_per_slide = cfg.getint("ppt", "landscape_per_slide", fallback=3)
        portrait_per_slide = cfg.getint("ppt", "portrait_per_slide", fallback=4)
        square_per_slide = cfg.getint("ppt", "square_per_slide", fallback=4)
        def _parse_ratio(value: str, default: float) -> float:
            if not value:
                return default
            v = value.strip()
            if "/" in v:
                a, b = v.split("/", 1)
                try:
                    return float(a) / float(b)
                except Exception:
                    return default
            try:
                return float(v)
            except Exception:
                return default
        portrait_threshold = _parse_ratio(cfg.get("ppt", "portrait_threshold", fallback="4/5"), 0.8)
        landscape_threshold = _parse_ratio(cfg.get("ppt", "landscape_threshold", fallback="16/9"), 16.0 / 9.0)
        title_font_size = cfg.getint("ppt", "title_font_size", fallback=28)
        content_font_size = cfg.getint("ppt", "content_font_size", fallback=10)
        title_left_in = cfg.getfloat("ppt", "title_left_in", fallback=0.5)
        title_top_in = cfg.getfloat("ppt", "title_top_in", fallback=0.3)
        title_width_in = cfg.getfloat("ppt", "title_width_in", fallback=12.4)
        title_height_in = cfg.getfloat("ppt", "title_height_in", fallback=1.25)
        content_left_in = cfg.getfloat("ppt", "content_left_in", fallback=0.5)
        content_top_in = cfg.getfloat("ppt", "content_top_in", fallback=0.8)
        content_width_in = cfg.getfloat("ppt", "content_width_in", fallback=13.0)
        content_height_in = cfg.getfloat("ppt", "content_height_in", fallback=6.5)
        title_font_name = cfg.get("ppt", "title_font_name", fallback="").strip()
        content_font_name = cfg.get("ppt", "content_font_name", fallback="").strip()
        template_path = cfg.get("ppt", "template_path", fallback="").strip()

        return CliConfig(
            recursive=recursive,
            formats=[ext.strip().lower() for ext in formats_csv.split(",") if ext.strip()],
            strategy=strategy,
            filename_separator=filename_separator,
            max_levels=max_levels,
            max_images_per_slide=max_images_per_slide,
            layout_strategy=layout_strategy,
            language=language,
            ppt_size=ppt_size,
            ppt_width_in=width_in,
            ppt_height_in=height_in,
            landscape_per_slide=landscape_per_slide,
            portrait_per_slide=portrait_per_slide,
            square_per_slide=square_per_slide,
            portrait_threshold=portrait_threshold,
            landscape_threshold=landscape_threshold,
            title_font_size=title_font_size,
            content_font_size=content_font_size,
            title_left_in=title_left_in,
            title_top_in=title_top_in,
            title_width_in=title_width_in,
            title_height_in=title_height_in,
            content_left_in=content_left_in,
            content_top_in=content_top_in,
            content_width_in=content_width_in,
            content_height_in=content_height_in,
            title_font_name=title_font_name,
            content_font_name=content_font_name,
            template_path=template_path,
        )

    def supported_extensions(self, override_formats: Optional[str]) -> List[str]:
        if override_formats:
            return [ext.strip().lower() for ext in override_formats.split(",") if ext.strip()]
        return list(self.formats or [])

    def slide_size_inches(self) -> tuple[float, float]:
        if self.ppt_width_in and self.ppt_height_in:
            return (self.ppt_width_in, self.ppt_height_in)
        key = (self.ppt_size or "16:9").lower()
        if key in {"16:9", "16x9", "widescreen"}:
            return (13.333, 7.5)
        if key in {"4:3", "4x3", "standard"}:
            return (10.0, 7.5)
        return (13.333, 7.5)


# -----------------------------
# Embedded image scanning/analyzer (from core/image_*.py)
# -----------------------------
class ImageScanner:
    def __init__(self, allowed_extensions: Iterable[str], recursive: bool = True) -> None:
        self.allowed = {ext.lower().lstrip(".") for ext in allowed_extensions}
        self.recursive = recursive

    def is_image(self, path: Path) -> bool:
        return path.is_file() and path.suffix.lower().lstrip(".") in self.allowed

    def scan(self, root: Path) -> List[Path]:
        if not root.exists():
            return []
        if root.is_file():
            return [root] if self.is_image(root) else []
        files: List[Path] = []
        iterator = root.rglob("*") if self.recursive else root.glob("*")
        for file in iterator:
            if self.is_image(file):
                files.append(file)
        files.sort()
        return files


@dataclass
class ImageInfo:
    path: str
    name: str
    width: int | None
    height: int | None
    orientation: str | None
    valid: bool


class ImageAnalyzer:
    def analyze(self, file: Path) -> Dict:
        try:
            with Image.open(file) as img:
                width, height = img.size
            if not width or not height or width <= 0 or height <= 0:
                return asdict(ImageInfo(str(file), file.name, None, None, None, False))
            if width > height:
                ori = "landscape"
            elif height > width:
                ori = "portrait"
            else:
                ori = "square"
            return asdict(ImageInfo(str(file), file.name, int(width), int(height), ori, True))
        except (UnidentifiedImageError, OSError, ValueError):
            return asdict(ImageInfo(str(file), file.name, None, None, None, False))

    def analyze_batch(self, files: Iterable[Path]) -> List[Dict]:
        return [self.analyze(f) for f in files]


# -----------------------------
# Embedded grouping (from core/grouper.py)
# -----------------------------
@dataclass
class GroupingOptions:
    strategy: str = "mixed"  # directory | filename | mixed
    filename_separator: str = "_"


GroupedBatch = Dict[str, object]


class GroupingStrategy:
    def __init__(self, options: Optional[GroupingOptions] = None) -> None:
        self.options = options or GroupingOptions()

    def _group_key_by_directory(self, img: Dict, root: Path) -> str:
        p = Path(str(img["path"]))
        try:
            rel = p.relative_to(root)
        except Exception:
            rel = p
        if rel.parent == Path(".") or str(rel.parent) == "":
            return "_root"
        return str(rel.parts[0])

    def _group_key_by_filename(self, img: Dict) -> str:
        name = str(img.get("name") or "")
        sep = self.options.filename_separator or "_"
        if sep in name:
            return name.split(sep, 1)[0]
        return name.rsplit(".", 1)[0]

    def group(self, analyzed_images: Iterable[Dict], root: Path) -> List[GroupedBatch]:
        groups: Dict[str, List[Dict]] = {}
        for img in analyzed_images:
            if not img.get("valid"):
                continue
            if self.options.strategy == "directory":
                key = self._group_key_by_directory(img, root)
            elif self.options.strategy == "filename":
                key = self._group_key_by_filename(img)
            else:
                k_dir = self._group_key_by_directory(img, root)
                key = k_dir if k_dir != "_root" else self._group_key_by_filename(img)
            groups.setdefault(key, []).append(img)
        ordered: List[GroupedBatch] = []
        for group_name in sorted(groups.keys()):
            images = sorted(groups[group_name], key=lambda x: str(x["name"]).lower())
            ordered.append({"group_name": group_name, "images": images})
        return ordered


# -----------------------------
# Embedded PPT generation (from core/ppt_generator.py)
# -----------------------------
@dataclass
class PPTOptions:
    max_images_per_slide: int = 4
    layout_strategy: str = "auto"
    language: str = "zh_CN"
    slide_width_in: float = 13.333
    slide_height_in: float = 7.5
    title_font_size: int = 28
    content_font_size: int = 10
    title_font_name: str = ""
    content_font_name: str = ""
    title_left_in: float = 0.5
    title_top_in: float = 0.3
    title_width_in: float = 12.4
    title_height_in: float = 1.25
    content_left_in: float = 0.5
    content_top_in: float = 0.8
    content_width_in: float = 13.0
    content_height_in: float = 6.5
    template_path: str = ""


class PPTGenerator:
    def __init__(self, options: Optional[PPTOptions] = None) -> None:
        self.options = options or PPTOptions()

    def _find_blank_layout(self, prs: Presentation) -> int:
        """Find blank layout (no placeholders) to avoid default font sizes."""
        try:
            for idx, layout in enumerate(prs.slide_layouts):
                try:
                    if len(layout.placeholders) == 0:
                        return idx
                except Exception:
                    continue
        except Exception:
            pass
        # Fallback to first layout
        return 0

    def _find_title_only_layout(self, prs: Presentation) -> int:
        """Return index of a 'Title Only' layout if available; otherwise fallback."""
        best_idx = None
        try:
            for idx, layout in enumerate(prs.slide_layouts):
                try:
                    has_title = False
                    for ph in layout.placeholders:
                        pf = getattr(ph, "placeholder_format", None)
                        if PP_PLACEHOLDER and pf and getattr(pf, "type", None) == PP_PLACEHOLDER.TITLE:
                            has_title = True
                            break
                        # Fallback heuristic: first placeholder with text_frame
                        if getattr(ph, "text_frame", None) is not None:
                            has_title = True
                            break
                    if has_title:
                        # prefer layouts with fewer placeholders (closer to title-only)
                        if best_idx is None or len(layout.placeholders) < len(prs.slide_layouts[best_idx].placeholders):
                            best_idx = idx
                except Exception:
                    continue
        except Exception:
            best_idx = None
        return best_idx if best_idx is not None else self._slide_layout_for(prs)

    def _remove_all_placeholders(self, slide) -> None:
        """Remove all placeholder shapes from the slide (including title and subtitle)."""
        to_delete = []
        for shape in list(slide.shapes):
            try:
                # Only delete placeholders (not regular shapes we added)
                if getattr(shape, "is_placeholder", False):
                    to_delete.append(shape)
            except Exception:
                continue
        # Delete safely
        for shape in to_delete:
            try:
                sp = shape._element  # lxml element
                sp.getparent().remove(sp)
            except Exception:
                continue

    def _remove_non_title_placeholders(self, slide, title_shape=None, title_placeholder_id=None) -> None:
        """Remove all placeholder shapes except the title one."""
        if title_shape is None:
            try:
                title_shape = slide.shapes.title
            except Exception:
                title_shape = None
        # Collect shapes to delete
        to_delete = []
        for shape in list(slide.shapes):
            try:
                # Skip title shape itself
                if shape is title_shape:
                    continue
                # Only delete placeholders (not regular shapes we added)
                if getattr(shape, "is_placeholder", False):
                    # Check if it's a title placeholder (should not delete)
                    try:
                        if hasattr(shape, 'placeholder_format'):
                            pf = shape.placeholder_format
                            # Check by type
                            if PP_PLACEHOLDER and hasattr(pf, 'type') and pf.type == PP_PLACEHOLDER.TITLE:
                                continue  # Don't delete title placeholder
                            # Also check by ID if provided
                            if title_placeholder_id is not None and hasattr(pf, 'idx') and pf.idx == title_placeholder_id:
                                continue  # Don't delete title placeholder
                    except Exception:
                        pass
                    to_delete.append(shape)
            except Exception:
                continue
        # Delete safely
        for shape in to_delete:
            try:
                sp = shape._element  # lxml element
                sp.getparent().remove(sp)
            except Exception:
                continue

    def _create_cover_page(self, prs: Presentation, slides: List[Dict[str, object]]) -> None:
        """Create cover page with group statistics and current date."""
        from collections import defaultdict
        from datetime import datetime
        
        # Extract first-level groups and count images
        group_counts: Dict[str, int] = defaultdict(int)
        total_images = 0
        
        for slide_spec in slides:
            breadcrumb = str(slide_spec.get("breadcrumb", ""))
            image_paths = slide_spec.get("images", [])
            image_count = len(image_paths) if isinstance(image_paths, list) else 0
            total_images += image_count
            
            # Extract first-level group from breadcrumb (format: "一级 / 二级 / 三级")
            if breadcrumb:
                first_level = breadcrumb.split(" / ")[0] if " / " in breadcrumb else breadcrumb
                group_counts[first_level] += image_count
        
        # Create cover slide
        layout_idx = self._find_blank_layout(prs)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        self._remove_all_placeholders(slide)
        
        # Layout parameters
        content_left = Inches(1.0)
        content_width = Inches(11.0)
        current_y = Inches(1.0)
        line_height = Inches(0.6)
        gap_between = Inches(0.3)
        
        # 1. Date textbox (independent)
        date_str = datetime.now().strftime("%Y年%m月%d日")
        date_shape = slide.shapes.add_textbox(
            content_left,
            current_y,
            content_width,
            line_height,
        )
        date_tf = date_shape.text_frame
        date_tf.clear()
        date_tf.word_wrap = True
        date_tf.margin_left = Inches(0.1)
        date_tf.margin_right = Inches(0.1)
        date_tf.margin_top = Inches(0.05)
        date_tf.margin_bottom = Inches(0.05)
        p = date_tf.paragraphs[0]
        run = p.add_run()
        run.text = f"日期: {date_str}"
        run.font.size = Pt(24)
        if self.options.content_font_name:
            run.font.name = self.options.content_font_name
        
        current_y += line_height + gap_between
        
        # 2. Groups textbox (with adaptive columns and font size)
        group_items = []
        if group_counts:
            for group_name, count in sorted(group_counts.items()):
                group_items.append(f"{group_name}: {count} 张")
        
        num_groups = len(group_items)
        groups_height = Inches(4.0)  # Available height for groups
        
        # Determine optimal layout: 1 column -> 2 columns -> 3 columns -> shrink font
        font_size = 24
        num_columns = 1
        
        # Try 1 column first
        estimated_height_1col = line_height * num_groups
        if estimated_height_1col <= groups_height:
            num_columns = 1
        else:
            # Try 2 columns
            items_per_col_2 = (num_groups + 1) // 2
            estimated_height_2col = line_height * items_per_col_2
            if estimated_height_2col <= groups_height:
                num_columns = 2
            else:
                # Try 3 columns
                items_per_col_3 = (num_groups + 2) // 3
                estimated_height_3col = line_height * items_per_col_3
                if estimated_height_3col <= groups_height:
                    num_columns = 3
                else:
                    # Need to shrink font
                    num_columns = 3
                    # Calculate required font size
                    items_per_col = items_per_col_3
                    required_height = groups_height
                    font_size = int(24 * (required_height / estimated_height_3col))
                    font_size = max(12, min(24, font_size))  # Clamp between 12 and 24
        
        # Calculate column layout
        gap = Inches(0.2)
        column_width = (content_width - gap * (num_columns - 1)) / num_columns
        items_per_column = (num_groups + num_columns - 1) // num_columns
        
        # Create groups textbox
        groups_shape = slide.shapes.add_textbox(
            content_left,
            current_y,
            content_width,
            groups_height,
        )
        groups_tf = groups_shape.text_frame
        groups_tf.clear()
        groups_tf.word_wrap = False
        groups_tf.margin_left = Inches(0.1)
        groups_tf.margin_right = Inches(0.1)
        groups_tf.margin_top = Inches(0.1)
        groups_tf.margin_bottom = Inches(0.1)
        
        # Build text with column layout
        # Each group should be on its own line, distributed across columns
        if num_columns == 1:
            # Single column: one item per line
            all_text = "\n".join(group_items)
        else:
            # Multiple columns: distribute items across columns
            lines = []
            for row in range(items_per_column):
                row_items = []
                for col in range(num_columns):
                    idx = col * items_per_column + row
                    if idx < num_groups:
                        row_items.append(group_items[idx])
                    else:
                        row_items.append("")
                # Join with tab for column alignment
                line = "\t".join(row_items)
                # Remove trailing tabs
                line = line.rstrip("\t")
                if line:
                    lines.append(line)
            all_text = "\n".join(lines)
        
        # Set up tab stops for columns
        p = groups_tf.paragraphs[0]
        p.clear()
        try:
            p.tab_stops.clear_all()
            for col in range(num_columns):
                tab_pos = content_left + col * (column_width + gap)
                p.tab_stops.add_tab_stop(tab_pos)
        except Exception:
            pass
        
        # Add text with calculated font size
        text_lines = all_text.split("\n")
        for i, line in enumerate(text_lines):
            if i > 0:
                p = groups_tf.add_paragraph()
                try:
                    for col in range(num_columns):
                        tab_pos = content_left + col * (column_width + gap)
                        p.tab_stops.add_tab_stop(tab_pos)
                except Exception:
                    pass
            run = p.add_run()
            run.text = line
            run.font.size = Pt(font_size)
            if self.options.content_font_name:
                run.font.name = self.options.content_font_name
        
        # Enable auto-fit as fallback
        try:
            if hasattr(groups_tf, 'auto_size'):
                from pptx.enum.text import MSO_AUTO_SIZE
                groups_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass
        
        current_y += groups_height + gap_between
        
        # 3. Total textbox (independent)
        total_shape = slide.shapes.add_textbox(
            content_left,
            current_y,
            content_width,
            line_height,
        )
        total_tf = total_shape.text_frame
        total_tf.clear()
        total_tf.word_wrap = True
        total_tf.margin_left = Inches(0.1)
        total_tf.margin_right = Inches(0.1)
        total_tf.margin_top = Inches(0.05)
        total_tf.margin_bottom = Inches(0.05)
        p = total_tf.paragraphs[0]
        run = p.add_run()
        run.text = f"总计: {total_images} 张"
        run.font.size = Pt(24)
        run.font.bold = True
        if self.options.content_font_name:
            run.font.name = self.options.content_font_name

    def generate_from_slides(self, slides: List[Dict[str, object]], output_path: Path) -> Path:
        prs = Presentation()
        prs.slide_width = Inches(self.options.slide_width_in)
        prs.slide_height = Inches(self.options.slide_height_in)
        
        # Create cover page with group statistics
        self._create_cover_page(prs, slides)
        
        for slide_spec in slides:
            title_text = str(slide_spec.get("title") or "")
            image_paths = [str(p) for p in (slide_spec.get("images") or [])]  # type: ignore[index]
            layout_kind = str(slide_spec.get("layout") or "")
            
            # Use blank layout to avoid placeholder default font sizes
            layout_idx = self._find_blank_layout(prs)
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            
            # Remove all placeholders (including "click to add title" and "click to add subtitle")
            self._remove_all_placeholders(slide)
            
            # Create title as textbox (not placeholder) to avoid default font size
            # Title position: top-left corner (left=0, top=0)
            title_shape = slide.shapes.add_textbox(
                Inches(0),
                Inches(0),
                Inches(self.options.title_width_in),
                Inches(self.options.title_height_in),
            )
            # Set title text and styling
            title_shape.text_frame.clear()
            p = title_shape.text_frame.paragraphs[0]
            p.clear()
            run = p.add_run()
            run.text = title_text
            # Title font size: half of configured size
            run.font.size = Pt(self.options.title_font_size / 2.0)
            if self.options.title_font_name:
                run.font.name = self.options.title_font_name
            # Left align title
            try:
                if PP_ALIGN:
                    p.alignment = PP_ALIGN.LEFT
            except Exception:
                pass
            
            # Load image dimensions
            images = []
            for p in image_paths:
                img_path = Path(p)
                width = None
                height = None
                try:
                    with Image.open(str(img_path)) as im:
                        width, height = im.size
                except Exception:
                    pass
                images.append({
                    "path": str(img_path),
                    "name": img_path.name,
                    "width": width,
                    "height": height
                })
            self._place_images_by_layout(slide, images, layout_kind)
        output_path = Path(output_path).resolve()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return output_path

    def _title_for_batch(self, group_name: str, batch_idx: int, batch_count: int) -> str:
        if batch_count <= 1:
            return group_name
        return f"{group_name} ({batch_idx}/{batch_count})"

    def _slide_layout_for(self, prs: Presentation) -> int:
        return 1 if len(prs.slide_layouts) > 1 else 0

    def _compute_slots(self, images: List[Dict]) -> int:
        if self.options.max_images_per_slide > 0:
            return self.options.max_images_per_slide
        return 4

    def _place_images(self, slide, images: List[Dict]) -> None:
        left_margin = Inches(0.5)
        top_margin = Inches(1.6)
        cell_size = Inches(4.0)
        gap = Inches(0.2)
        positions: List[Tuple[float, float]] = []
        for r in range(2):
            for c in range(2):
                x = left_margin + c * (cell_size + gap)
                y = top_margin + r * (cell_size + gap)
                positions.append((x, y))
        for idx, img in enumerate(images[:4]):
            x, y = positions[idx]
            w_px = img.get("width")
            h_px = img.get("height")
            if not (w_px and h_px):
                try:
                    with Image.open(str(img["path"])) as im:
                        w_px, h_px = im.size
                except Exception:
                    w_px, h_px = (1000, 1000)
            w_px = float(w_px)
            h_px = float(h_px)
            if w_px >= h_px:
                disp_w = cell_size
                disp_h = cell_size * (h_px / w_px)
            else:
                disp_h = cell_size
                disp_w = cell_size * (w_px / h_px)
            offset_x = x + (cell_size - disp_w) / 2.0
            offset_y = y + (cell_size - disp_h) / 2.0
            slide.shapes.add_picture(img["path"], offset_x, offset_y, width=disp_w, height=disp_h)
            name = str(img.get("name") or "")
            w = img.get("width")
            h = img.get("height")
            # Format: "图片名 - (wxh)"
            caption = f"{name} - ({w}x{h})" if (w and h) else name
            tb = slide.shapes.add_textbox(x, y + cell_size - Inches(0.3), width=cell_size, height=Inches(0.3))
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            try:
                if PP_ALIGN:
                    p.alignment = PP_ALIGN.CENTER
            except Exception:
                pass
            run = p.add_run()
            run.text = caption
            run.font.size = Pt(self.options.content_font_size)
            if self.options.content_font_name:
                run.font.name = self.options.content_font_name

    def _place_images_by_layout(self, slide, images: List[Dict], layout_kind: str) -> None:
        lk = layout_kind.lower()
        # Use configured slide width for centering
        slide_width = Inches(self.options.slide_width_in)
        if lk == "landscape":
            self._place_landscape_vertical(slide, images, per_page=self.options.max_images_per_slide or 3, slide_width=slide_width)
        elif lk == "portrait":
            self._place_portrait_horizontal(slide, images, per_page=self.options.max_images_per_slide or 4, slide_width=slide_width)
        else:
            self._place_square_grid(slide, images, per_page=self.options.max_images_per_slide or 4, slide_width=slide_width)

    def _place_square_grid(self, slide, images: List[Dict], per_page: int = 4, slide_width=None) -> None:
        # Center content area horizontally on slide
        if slide_width is None:
            slide_width = Inches(self.options.slide_width_in)
        content_width = Inches(self.options.content_width_in)
        centered_left = (slide_width - content_width) / 2.0
        left = centered_left
        top = Inches(self.options.content_top_in)
        width = content_width
        height = Inches(self.options.content_height_in)
        cols = 2
        rows = 2
        gap = Inches(0.2)
        cell_w = (width - gap) / cols
        cell_h = (height - gap) / rows

        positions: List[Tuple[float, float]] = []
        for r in range(rows):
            for c in range(cols):
                x = left + c * (cell_w + (gap if c < cols else 0))
                y = top + r * (cell_h + (gap if r < rows else 0))
                positions.append((x, y))

        for idx, img in enumerate(images[:4]):
            x, y = positions[idx]
            w_px = img.get("width")
            h_px = img.get("height")
            if not (w_px and h_px):
                try:
                    with Image.open(str(img["path"])) as im:
                        w_px, h_px = im.size
                except Exception:
                    w_px, h_px = (1000, 1000)
            w_px = float(w_px)
            h_px = float(h_px)
            # fit within cell preserving aspect, reserve caption space
            max_h = cell_h - Inches(0.3)
            if w_px >= h_px:
                disp_w = cell_w
                disp_h = disp_w * (h_px / w_px)
                if disp_h > max_h:
                    disp_h = max_h
                    disp_w = disp_h * (w_px / h_px)
            else:
                disp_h = max_h
                disp_w = disp_h * (w_px / h_px)
                if disp_w > cell_w:
                    disp_w = cell_w
                    disp_h = disp_w * (h_px / w_px)
            offset_x = x + (cell_w - disp_w) / 2.0
            offset_y = y + (max_h - disp_h) / 2.0
            slide.shapes.add_picture(img["path"], offset_x, offset_y, width=disp_w, height=disp_h)
            # caption
            name = str(img.get("name") or "")
            w = img.get("width")
            h = img.get("height")
            # Format: "图片名 - (wxh)"
            caption = f"{name} - ({w}x{h})" if (w and h) else name
            tb = slide.shapes.add_textbox(x, y + max_h, width=cell_w, height=Inches(0.3))
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            try:
                if PP_ALIGN:
                    p.alignment = PP_ALIGN.CENTER
            except Exception:
                pass
            run = p.add_run()
            run.text = caption
            run.font.size = Pt(self.options.content_font_size)
            if self.options.content_font_name:
                run.font.name = self.options.content_font_name

    def _place_landscape_vertical(self, slide, images: List[Dict], per_page: int = 3, slide_width=None) -> None:
        # Center content area horizontally on slide
        if slide_width is None:
            slide_width = Inches(self.options.slide_width_in)
        content_width = Inches(self.options.content_width_in)
        centered_left = (slide_width - content_width) / 2.0
        left = centered_left
        top = Inches(self.options.content_top_in)
        width = content_width
        height = Inches(self.options.content_height_in)
        cell_h = height / float(max(1, min(per_page, len(images))))
        cell_w = width
        gap = Inches(0.2)
        y = top
        for idx, img in enumerate(images[:per_page]):
            # compute placement while preserving aspect ratio into area (cell_w x cell_h)
            w_px = img.get("width")
            h_px = img.get("height")
            if not (w_px and h_px):
                try:
                    with Image.open(str(img["path"])) as im:
                        w_px, h_px = im.size
                except Exception:
                    w_px, h_px = (1000, 600)
            w_px = float(w_px)
            h_px = float(h_px)
            # caption beneath: 1 line (name | w x h), reserve 0.3in
            cap_h = Inches(0.3)
            # fit to width primarily, accounting for caption space
            disp_w = cell_w
            disp_h = disp_w * (h_px / w_px)
            # Adjust image height if overlapping with caption
            if disp_h > cell_h - cap_h:
                disp_h = max(Inches(0.5), cell_h - cap_h)
                disp_w = disp_h * (w_px / h_px)
            # Also check against gap
            if disp_h > cell_h - gap:
                disp_h = cell_h - gap
                disp_w = disp_h * (w_px / h_px)
            x = left + (cell_w - disp_w) / 2.0
            # Add picture only once with final dimensions
            slide.shapes.add_picture(img["path"], x, y, width=disp_w, height=disp_h)
            name = str(img.get("name") or "")
            w = img.get("width")
            h = img.get("height")
            # Format: "图片名 - (wxh)"
            caption = f"{name} - ({w}x{h})" if (w and h) else name
            tb = slide.shapes.add_textbox(left, y + disp_h, width=cell_w, height=cap_h)
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            if PP_ALIGN:
                p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = caption
            r.font.size = Pt(self.options.content_font_size)
            if self.options.content_font_name:
                r.font.name = self.options.content_font_name
            y += cell_h

    def _place_portrait_horizontal(self, slide, images: List[Dict], per_page: int = 4, slide_width=None) -> None:
        # Center content area horizontally on slide
        if slide_width is None:
            slide_width = Inches(self.options.slide_width_in)
        content_width = Inches(self.options.content_width_in)
        centered_left = (slide_width - content_width) / 2.0
        left = centered_left
        top = Inches(self.options.content_top_in)
        width = content_width
        height = Inches(self.options.content_height_in)
        gap = Inches(0.3)  # Gap between cells
        cap_h = Inches(0.3)  # caption height
        # Calculate cell width with gaps between cells
        num_cells = max(1, min(per_page, len(images)))
        total_gaps = gap * (num_cells - 1) if num_cells > 1 else Inches(0)
        available_width = width - total_gaps
        cell_w = available_width / float(num_cells)
        cell_h = height
        x = left
        for idx, img in enumerate(images[:per_page]):
            w_px = img.get("width")
            h_px = img.get("height")
            if not (w_px and h_px):
                try:
                    with Image.open(str(img["path"])) as im:
                        w_px, h_px = im.size
                except Exception:
                    w_px, h_px = (800, 1200)
            w_px = float(w_px)
            h_px = float(h_px)
            # fit to height primarily (excluding caption area)
            caption_height = Inches(0.4)  # Match caption height used below
            img_area_h = cell_h - caption_height
            disp_h = img_area_h
            disp_w = disp_h * (w_px / h_px)
            # Leave some margin in cell for image
            img_gap = Inches(0.1)
            if disp_w > cell_w - img_gap:
                disp_w = cell_w - img_gap
                disp_h = disp_w * (h_px / w_px)
            # Center image horizontally in its cell
            img_x = x + (cell_w - disp_w) / 2.0
            # Center image vertically in image area (above caption)
            img_y = top + (img_area_h - disp_h) / 2.0
            slide.shapes.add_picture(img["path"], img_x, img_y, width=disp_w, height=disp_h)
            # caption at bottom of cell: 1 line (name - wxh)
            name = str(img.get("name") or "")
            w = img.get("width")
            h = img.get("height")
            # Format: "图片名 - (wxh)"
            caption = f"{name} - ({w}x{h})" if (w and h) else name
            # Textbox at bottom of cell, fixed width with word wrap
            # Increase height slightly to accommodate potential wrapping
            caption_height = Inches(0.4)  # Slightly taller to allow wrapping
            tb = slide.shapes.add_textbox(x, top + cell_h - caption_height, width=cell_w, height=caption_height)
            tf = tb.text_frame
            tf.clear()
            # Enable word wrap for long text
            tf.word_wrap = True
            # Set fixed margins to ensure text doesn't touch edges
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.02)
            tf.margin_bottom = Inches(0.02)
            p = tf.paragraphs[0]
            # Center text horizontally
            if PP_ALIGN:
                p.alignment = PP_ALIGN.CENTER
            # Align text to top of textbox (vertical alignment)
            try:
                from pptx.enum.text import MSO_ANCHOR
                tf.vertical_anchor = MSO_ANCHOR.TOP
            except Exception:
                pass
            r = p.add_run()
            r.text = caption
            r.font.size = Pt(self.options.content_font_size)
            if self.options.content_font_name:
                r.font.name = self.options.content_font_name
            # Move to next cell position (with gap)
            x += cell_w + gap

    def generate(self, groups: List[Dict], output_path: Path) -> Path:
        prs = Presentation()
        prs.slide_width = Inches(self.options.slide_width_in)
        prs.slide_height = Inches(self.options.slide_height_in)
        for group in groups:
            group_name = str(group["group_name"])
            images: List[Dict] = list(group["images"])
            if not images:
                continue
            max_per = self._compute_slots(images)
            batches = [images[i : i + max_per] for i in range(0, len(images), max_per)]
            for idx, batch in enumerate(batches, start=1):
                layout_idx = self._slide_layout_for(prs)
                slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
                title = self._title_for_batch(group_name, idx, len(batches))
                # Set title with proper font size
                try:
                    title_shape = slide.shapes.title
                    if title_shape is None:
                        title_shape = slide.shapes.add_textbox(
                            Inches(self.options.title_left_in),
                            Inches(self.options.title_top_in),
                            Inches(self.options.title_width_in),
                            Inches(self.options.title_height_in),
                        )
                    # If using placeholder, we need to clear it properly
                    if hasattr(title_shape, 'placeholder') and title_shape.placeholder:
                        title_shape.text_frame.clear()
                    else:
                        title_shape.text_frame.clear()
                    
                    p = title_shape.text_frame.paragraphs[0]
                    # Clear all runs first
                    p.clear()
                    run = p.add_run()
                    run.text = title
                    # Set font size
                    run.font.size = Pt(self.options.title_font_size)
                    # Force set font size via XML to ensure it's applied
                    try:
                        from pptx.oxml.ns import qn
                        rPr = run._r.get_or_add_rPr()
                        sz = rPr.find(qn('w:sz'))
                        if sz is None:
                            from pptx.oxml import parse_xml
                            sz = parse_xml(f'<w:sz xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" w:val="{int(self.options.title_font_size * 2)}"/>')
                            rPr.append(sz)
                        else:
                            sz.set(qn('w:val'), str(int(self.options.title_font_size * 2)))
                    except Exception:
                        pass
                    if self.options.title_font_name:
                        run.font.name = self.options.title_font_name
                    
                    try:
                        if PP_ALIGN:
                            p.alignment = PP_ALIGN.LEFT
                    except Exception:
                        pass
                except Exception:
                    # Fallback to simple text assignment
                    try:
                        slide.shapes.title.text = title
                    except Exception:
                        pass
                self._place_images(slide, batch)
        output_path = Path(output_path).resolve()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return output_path


# -----------------------------
# Embedded structure MD (from core/structure_md.py)
# -----------------------------
HEADER = "# i2pptt structure"


def _title_for_batch_md(group_name: str, idx: int, total: int) -> str:
    return f"{group_name} ({idx}/{total})" if total > 1 else group_name


def _filename_tokens(stem: str, sep: str) -> List[str]:
    # split by custom separator first, then by non-alnum unicode boundaries
    if sep and sep in stem:
        primary = stem.split(sep)
    else:
        primary = [stem]
    tokens: List[str] = []
    for chunk in primary:
        # Split on hyphen, whitespace, underscore, dot
        parts = re.split(r"[-\s_\.]+", chunk)
        for p in parts:
            if p:
                tokens.append(p)
    return tokens


def _clean_heading(name: str) -> str:
    # strip leading/trailing spaces, hyphens, underscores; keep middle intact
    return re.sub(r'^[\s\-_]+|[\s\-_]+$', '', name)

def _rotate_if_exists(path: Path) -> None:
    """If the target file exists, rename it to filename[NN].ext (NN starts at 01)."""
    path = Path(path)
    if not path.exists():
        return
    parent = path.parent
    name = path.stem
    suffix = path.suffix  # includes the dot
    # If name already has [NN], we still create a new [NN+1] to avoid overwriting
    index = 1
    while True:
        rotated = parent / f"{name}[{index:02d}]{suffix}"
        if not rotated.exists():
            try:
                path.rename(rotated)
            except OSError:
                # As a fallback, try a higher index next loop in case of race
                index += 1
                continue
            break
        index += 1


def write_structure_md_hier(analyzed_images: List[Dict], root: Path, max_levels: int, filename_sep: str, output_md: Path, ini_header: List[str] | None = None) -> Path:
    lines: List[str] = []
    lines.append(HEADER)
    lines.append("")
    lines.append(f"root: {str(root)}")
    lines.append("")
    if ini_header:
        lines.append("ini:")
        for entry in ini_header:
            lines.append(f"- {entry}")
        lines.append("")
    # Build a flat list of items with rel parts and filename tokens
    items: List[Dict] = []
    for img in analyzed_images:
        if not img.get("valid"):
            continue
        p = Path(str(img["path"]))
        try:
            rel = p.relative_to(root)
        except Exception:
            rel = p
        dirs = list(rel.parts[:-1])
        stem = Path(rel.name).stem
        tokens = _filename_tokens(stem, filename_sep)
        items.append({"img": img, "rel": rel, "dirs": dirs, "tokens": tokens})

    # Sort items by rel path for stable output
    items.sort(key=lambda x: str(x["rel"]))

    # Precompute candidate levels (name, source) up to max_levels for each item
    # and counts for prefixes; token-based headings with count==1 will be skipped.
    from collections import defaultdict
    prefix_counts: Dict[tuple, int] = defaultdict(int)
    full_levels_per_item: List[List[tuple[str, str]]] = []
    for it in items:
        dirs: List[str] = it["dirs"]
        tokens: List[str] = it["tokens"]
        levels: List[tuple[str, str]] = []
        dir_idx = 0
        tok_idx = 0
        for level in range(max_levels):
            if dir_idx < len(dirs):
                levels.append((dirs[dir_idx], "dir"))
                dir_idx += 1
            elif tok_idx < len(tokens):
                levels.append((tokens[tok_idx], "token"))
                tok_idx += 1
            else:
                break
        # record prefix counts for all prefixes
        prefix: List[str] = []
        for i, (name, _src) in enumerate(levels, start=1):
            prefix.append(name)
            prefix_counts[(i, tuple(prefix))] += 1
        full_levels_per_item.append(levels)

    # Emit headings and items by hierarchical grouping up to max_levels
    # We will keep track of last emitted path of headings to only write changed ones
    last_headings: List[str] = []

    # Compute LCP of stems under each parent directory prefix (up to max_levels-1 depth)
    stems_by_parent: Dict[tuple[str, ...], List[str]] = defaultdict(list)
    for it in items:
        dirs: List[str] = it["dirs"]
        parent_depth = min(len(dirs), max_levels - 1)
        parent_key = tuple(dirs[:parent_depth])
        if parent_depth >= 1:
            stems_by_parent[parent_key].append(Path(it["rel"]).stem)

    def longest_common_prefix(strings: List[str]) -> str:
        if not strings:
            return ""
        s1 = min(strings)
        s2 = max(strings)
        for i, c in enumerate(s1):
            if i >= len(s2) or c != s2[i]:
                return s1[:i]
        return s1

    lcp_by_parent: Dict[tuple[str, ...], str] = {
        key: longest_common_prefix(stems).strip() for key, stems in stems_by_parent.items() if stems
    }

    for it, levels in zip(items, full_levels_per_item):
        # Derive effective names: keep directory levels; then use LCP for next level if available,
        # otherwise keep token level only if count > 1
        effective_names: List[str] = []
        # 1) append directory levels
        for (name, src) in levels:
            if src == "dir":
                effective_names.append(name)
                if len(effective_names) == max_levels:
                    break
            else:
                break
        # 2) next level from LCP (if any and room)
        if len(effective_names) < max_levels:
            parent_key = tuple(effective_names)
            lcp = lcp_by_parent.get(parent_key, "")
            if lcp:
                effective_names.append(lcp)
            else:
                # fallback to token rule based on counts
                i_level = len(effective_names) + 1
                prefix_names = tuple(name for (name, _src) in levels[:i_level])
                if prefix_counts.get((i_level, prefix_names), 0) > 1:
                    # find first token after dir levels
                    for (name, src) in levels[len(effective_names):]:
                        if src == "token":
                            effective_names.append(name)
                            break

        # Clean headings by removing leading/trailing separators
        cleaned_names = [_clean_heading(n) for n in effective_names if _clean_heading(n)]
        # Write headings for the common prefix difference with last_headings
        common = 0
        while common < len(cleaned_names) and common < len(last_headings) and cleaned_names[common] == last_headings[common]:
            common += 1
        # Close/open headings by simply writing from the divergence point
        for level in range(common, len(cleaned_names)):
            depth = level + 1  # 1-based
            prefix = "#" * depth
       	# ensure non-empty after cleaning
            lines.append(f"{prefix} {cleaned_names[level]}")
        last_headings = cleaned_names

        # Emit the item entry under current headings
        img = it["img"]
        rel = it["rel"]
        width = img.get("width")
        height = img.get("height")
        name = img.get("name") or Path(rel).name
        meta = f"{width}x{height}" if (width and height) else ""
        orientation = (img.get("orientation") or "").strip()
        lines.append(f"- path: {rel}")
        lines.append(f"  filename: {Path(rel).name}")
        lines.append(f"  title: {name}")
        lines.append(f"  size: {meta}")
        lines.append(f"  orientation: {orientation}")
    lines.append("")
    output_md = output_md.resolve()
    output_md.parent.mkdir(parents=True, exist_ok=True)
    # Rotate existing MD if present
    _rotate_if_exists(output_md)
    output_md.write_text("\n".join(lines), encoding="utf-8")
    return output_md


def parse_structure_md(md_path: Path) -> List[Dict[str, object]]:
    text = Path(md_path).read_text(encoding="utf-8")
    lines = [line.rstrip() for line in text.splitlines()]
    slides: List[Dict[str, object]] = []
    # Track heading stack for breadcrumb (e.g., # A, ## B, ### C -> A / B / C)
    heading_stack: List[tuple[int, str]] = []
    current: Dict[str, object] | None = None
    i = 0
    while i < len(lines):
        line = lines[i]
        if line.startswith("#"):
            # any markdown heading level starts a new section
            hashes = len(line) - len(line.lstrip("#"))
            title = line.lstrip("#").strip()
            # update heading stack to this level
            while heading_stack and heading_stack[-1][0] >= hashes:
                heading_stack.pop()
            heading_stack.append((hashes, title))
            # compute breadcrumb from stack (order by level)
            breadcrumb = " / ".join(name for (_lvl, name) in heading_stack)
            if current and current.get("images"):
                slides.append(current)
            current = {"title": title, "breadcrumb": breadcrumb, "images": []}
            i += 1
            continue
        if line.startswith("- path:"):
            if current is None:
                i += 1
                continue
            path_value = line.split(":", 1)[1].strip()
            title_value = ""
            size_value = ""
            # read following indented attributes
            j = i + 1
            while j < len(lines) and lines[j].startswith("  "):
                sub = lines[j].strip()
                if sub.startswith("title:"):
                    title_value = sub.split(":", 1)[1].strip()
                elif sub.startswith("size:"):
                    size_value = sub.split(":", 1)[1].strip()
                j += 1
            imgs = current.get("images")
            if isinstance(imgs, list):
                imgs.append(path_value)
            i = j
            continue
        i += 1
    if current and current.get("images"):
        slides.append(current)
    return slides


@click.group(context_settings={"help_option_names": ["-h", "--help"]})
@click.version_option("0.1.0", prog_name="i2pptt")
def cli() -> None:
    """i2pptt command line interface."""


def _resolve_path(path: str) -> Path:
    p = Path(path).expanduser()
    if not p.exists():
        raise click.ClickException(f"Path not found: {p}")
    return p


@cli.command("scan")
@click.option("--directory", "--dir", "-d", required=True, type=str, help="Root directory of images.")
@click.option("--filename", "-f", required=True, type=str, help="Target PPT filename (used to derive <filename>_structure.md).")
@click.option("--recursive/--no-recursive", default=True, help="Scan subdirectories recursively.")
@click.option("--formats", default=None, help="Comma-separated extensions, e.g. png,jpg,jpeg,webp.")
@click.option("--config", "config_path", type=str, default=None, help="Path to i2pptt.ini.")
def scan(directory: str, filename: str, recursive: bool, formats: Optional[str], config_path: Optional[str]) -> None:
    """Scan directory and produce structure markdown at <filename>_structure.md."""
    cfg = CliConfig.load(config_path)
    allowed = cfg.supported_extensions(formats)
    root = _resolve_path(directory)

    scanner = ImageScanner(allowed_extensions=allowed, recursive=recursive)
    analyzer = ImageAnalyzer()
    files = scanner.scan(root)
    analyzed = analyzer.analyze_batch(files)

    base = Path(filename)
    if base.suffix.lower() == ".pptx":
        base = base.with_suffix("")
    md_path = base.with_name(f"{base.name}_structure.md")
    # Build ini header summary
    ini_used = config_path if config_path else "cli/i2pptt.ini"
    ini_header = [
        f"path = {ini_used}",
        f"ppt.size = {cfg.ppt_size}",
        f"ppt.slide_inches = {cfg.slide_size_inches()[0]}x{cfg.slide_size_inches()[1]}",
        f"ppt.landscape_per_slide = {cfg.landscape_per_slide}",
        f"ppt.portrait_per_slide = {cfg.portrait_per_slide}",
        f"ppt.square_per_slide = {cfg.square_per_slide}",
        f"ppt.portrait_threshold = {cfg.portrait_threshold}",
        f"ppt.landscape_threshold = {cfg.landscape_threshold}",
        f"title.font = {cfg.title_font_name or '-'} {cfg.title_font_size}pt",
        f"title.box_in = left {cfg.title_left_in}, top {cfg.title_top_in}, width {cfg.title_width_in}, height {cfg.title_height_in}",
        f"content.font = {cfg.content_font_name or '-'} {cfg.content_font_size}pt",
        f"content.box_in = left {cfg.content_left_in}, top {cfg.content_top_in}, width {cfg.content_width_in}, height {cfg.content_height_in}",
    ]
    write_structure_md_hier(
        analyzed,
        root=root,
        max_levels=cfg.max_levels,
        filename_sep=cfg.filename_separator,
        output_md=md_path,
        ini_header=ini_header,
    )
    click.echo(str(md_path))


@cli.command("merge")
@click.option("--directory", "--dir", "-d", required=True, type=str, help="Root directory of images.")
@click.option("--filename", "-f", required=True, type=str, help="Target PPT filename. Reads <filename>_structure.md.")
@click.option("--config", "config_path", type=str, default=None, help="Path to i2pptt.ini.")
def merge(directory: str, filename: str, config_path: Optional[str]) -> None:
    """Generate PPT from structure markdown."""
    cfg = CliConfig.load(config_path)
    root = _resolve_path(directory)

    base = Path(filename)
    out_ppt = base if base.suffix.lower() == ".pptx" else base.with_suffix(".pptx")
    md_path = base.with_suffix("").with_name(f"{base.with_suffix('').name}_structure.md")

    if not md_path.exists():
        raise click.ClickException(f"Structure file not found: {md_path}")

    raw_sections = parse_structure_md(md_path)
    # Build slides based on orientation groups and pagination
    slides: List[Dict[str, object]] = []
    def classify(img_path: str) -> str:
        try:
            with Image.open(img_path) as im:
                w, h = im.size
        except Exception:
            w, h = (1000, 1000)
        ratio = float(w) / float(h) if h else 1.0
        if ratio < cfg.portrait_threshold:
            return "portrait"
        if ratio > cfg.landscape_threshold:
            return "landscape"
        return "square"
    for sect in raw_sections:
        breadcrumb = str(sect.get("breadcrumb") or sect.get("title") or "")
        imgs = list(sect.get("images") or [])
        abs_imgs = [str((root / p).resolve()) for p in imgs]
        # orientation grouping
        groups: Dict[str, List[str]] = {"landscape": [], "portrait": [], "square": []}
        for p in abs_imgs:
            g = classify(p)
            groups[g].append(p)
        # paginate each group
        for kind in ("landscape", "portrait", "square"):
            arr = groups[kind]
            if not arr:
                continue
            if kind == "landscape":
                per = max(1, int(cfg.landscape_per_slide))
            elif kind == "portrait":
                per = max(1, int(cfg.portrait_per_slide))
            else:
                per = max(1, int(cfg.square_per_slide))
            batches = [arr[i : i + per] for i in range(0, len(arr), per)]
            for idx, batch in enumerate(batches, start=1):
                slides.append({
                    "title": f"{breadcrumb} ({idx}/{len(batches)})",
                    "breadcrumb": breadcrumb,  # Preserve breadcrumb for cover page statistics
                    "images": batch,
                    "layout": kind,
                })

    # Resolve template path relative to config file if needed
    template_path = cfg.template_path
    if template_path:
        template_path_obj = Path(template_path).expanduser()
        if not template_path_obj.is_absolute():
            # If relative, try relative to config file location
            if config_path:
                config_dir = Path(config_path).parent
                template_path_obj = (config_dir / template_path).resolve()
            else:
                template_path_obj = Path(template_path).resolve()
        template_path = str(template_path_obj) if template_path_obj.exists() else ""
    
    generator = PPTGenerator(PPTOptions(
        max_images_per_slide=cfg.max_images_per_slide,
        layout_strategy=cfg.layout_strategy,
        language=cfg.language,
        slide_width_in=cfg.slide_size_inches()[0],
        slide_height_in=cfg.slide_size_inches()[1],
        title_font_size=cfg.title_font_size,
        content_font_size=cfg.content_font_size,
        title_left_in=cfg.title_left_in,
        title_top_in=cfg.title_top_in,
        title_width_in=cfg.title_width_in,
        title_height_in=cfg.title_height_in,
        content_left_in=cfg.content_left_in,
        content_top_in=cfg.content_top_in,
        content_width_in=cfg.content_width_in,
        content_height_in=cfg.content_height_in,
        title_font_name=cfg.title_font_name,
        content_font_name=cfg.content_font_name,
        template_path=template_path,
    ))
    # Rotate existing PPTX if present before writing
    _rotate_if_exists(out_ppt)
    ppt_path = generator.generate_from_slides(slides, output_path=out_ppt)
    click.echo(str(ppt_path))


# Deprecated commands 'group' and 'generate' are removed in favor of 'scan' and 'merge'.


if __name__ == "__main__":
    try:
        cli()
    except click.ClickException as exc:
        click.echo(f"Error: {exc}", err=True)
        sys.exit(2)
    except KeyboardInterrupt:
        click.echo("Aborted.", err=True)
        sys.exit(130)


